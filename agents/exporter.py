
import json
import os
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

EXPORTS_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "exports")

def export_to_csv(deals, filename="newsletter.csv"):
    if not os.path.exists(EXPORTS_DIR):
        os.makedirs(EXPORTS_DIR)
    filepath = os.path.join(EXPORTS_DIR, filename)
    df = pd.DataFrame(deals)
    df.to_csv(filepath, index=False)
    return filepath

def export_to_json(deals, filename="newsletter.json"):
    filepath = os.path.join(EXPORTS_DIR, filename)
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(deals, f, ensure_ascii=False, indent=2)
    return filepath

def export_to_excel(deals, filename="newsletter.xlsx"):
    filepath = os.path.join(EXPORTS_DIR, filename)
    df = pd.DataFrame(deals)
    df.to_excel(filepath, index=False)
    return filepath

def export_to_word(newsletter_text, deals, filename="newsletter.docx"):
    doc = Document()
    doc.add_heading("FMCG Intelligence Newsletter", 0)
    
    for paragraph in newsletter_text.split("\n\n"):
        if paragraph.strip():
            doc.add_paragraph(paragraph)
            
    doc.save(os.path.join(EXPORTS_DIR, filename))
    return os.path.join(EXPORTS_DIR, filename)

def export_to_powerpoint(newsletter_text, deals, filename="newsletter.pptx"):
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "FMCG Intelligence Newsletter"
    
    paragraphs = newsletter_text.split("\n\n")
    for para in paragraphs[:10]:
        if para.strip():
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = para.split("\n")[0]
            slide.placeholders[1].text = para
            
    prs.save(os.path.join(EXPORTS_DIR, filename))
    return os.path.join(EXPORTS_DIR, filename)

if __name__ == "__main__":
    from agents.extraction import load_deals
    from agents.newsletter import generate_newsletter
    deals = load_deals()
    newsletter = generate_newsletter(deals)
    
    export_to_csv(deals)
    export_to_json(deals)
    export_to_excel(deals)
    export_to_word(newsletter, deals)
    export_to_powerpoint(newsletter, deals)

